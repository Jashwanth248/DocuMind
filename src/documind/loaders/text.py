from pathlib import Path
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from ..schemas import ContentUnit

TEXT_SUFFIXES = {".txt", ".md", ".csv", ".json", ".xml", ".py", ".sql", ".rtf"}


def load_textual(path: Path, source_id: str) -> list[ContentUnit]:
    suffix = path.suffix.lower()
    if suffix in TEXT_SUFFIXES:
        text = path.read_text(encoding="utf-8", errors="ignore")
        return [ContentUnit(source_id, path.name, str(path), "text", text)]
    if suffix in {".html", ".htm"}:
        soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "html.parser")
        return [ContentUnit(source_id, path.name, str(path), "text", soup.get_text(" ", strip=True))]
    if suffix == ".docx":
        doc = Document(path)
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return [ContentUnit(source_id, path.name, str(path), "document", text)]
    if suffix == ".pptx":
        prs = Presentation(path)
        units = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if parts:
                units.append(ContentUnit(source_id, path.name, str(path), "presentation", "\n".join(parts), f"slide {i}"))
        return units
    if suffix == ".xlsx":
        wb = load_workbook(path, read_only=True, data_only=True)
        units = []
        for ws in wb.worksheets:
            rows = []
            for row in ws.iter_rows(values_only=True):
                vals = ["" if v is None else str(v) for v in row]
                if any(vals):
                    rows.append(" | ".join(vals))
            if rows:
                units.append(ContentUnit(source_id, path.name, str(path), "spreadsheet", "\n".join(rows), f"sheet {ws.title}"))
        return units
    return []
